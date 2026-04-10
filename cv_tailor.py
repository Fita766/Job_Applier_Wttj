import os
import time
import json
import re
from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.enum.dml import MSO_COLOR_TYPE
from pptx.dml.color import RGBColor
import datetime
import config
try:
    import comtypes.client
except ImportError:
    pass

from ai_helper import _call_openai_with_mistral_fallback
import config

def extract_experience_shapes(pptx_path):
    """Extrait les blocs de texte susceptibles d'etre des descriptions d'experiences."""
    prs = Presentation(pptx_path)
    shapes_data = []
    
    for slide_idx, slide in enumerate(prs.slides):
        for shape_idx, shape in enumerate(slide.shapes):
            if not shape.has_text_frame:
                continue
            text = shape.text.strip()
            # Heuristique: Un bloc d'experience a generalement plus de 40 caracteres
            if len(text) > 40 and "@" not in text and "Tel" not in text and "http" not in text:
                # On "nettoie" le texte des puces manuelles pour le JSON envoye a l'IA
                # pour eviter que l'IA ne les reinclue dans son texte (le script les rajoute)
                clean_text = re.sub(r"^[->*\->*]\s*", "", text, flags=re.MULTILINE)
                char_count = len(text)
                
                # V13 : On compte le nombre exact de paragraphes/puces d'origine
                num_orig_paragraphs = len([p for p in shape.text_frame.paragraphs if p.text.strip()])
                # V17 : Forcage extreme selon les instructions (5 pour Exp 1/2, 4 pour Exp 3)
                if "AXA" in text:
                    target_paragraphs = 5
                elif "Helpline" in text:
                    target_paragraphs = 5
                elif "TSN" in text:
                    target_paragraphs = 3
                else:
                    target_paragraphs = max(3, num_orig_paragraphs - 1)
                
                shapes_data.append({
                    "slide_idx": slide_idx,
                    "shape_idx": shape_idx,
                    "text": clean_text,
                    "char_count": char_count,
                    "min_chars": int(char_count * 0.90),
                    "max_chars": char_count, # V17 : Retour a 100% pour eviter de tronquer les belles phrases
                    "exact_bullet_count": target_paragraphs
                })
    return shapes_data

def update_pptx_shapes(pptx_path, updates, output_path):
    """Met a jour les shapes du PPTX avec les nouveaux textes."""
    prs = Presentation(pptx_path)
    
    # updates est une liste de {"slide_idx": 0, "shape_idx": 5, "text": "nouveau texte"}
    for update in updates:
        slide_idx = update.get("slide_idx")
        shape_idx = update.get("shape_idx")
        new_text = update.get("text", "").strip()
        
        if slide_idx is None or shape_idx is None or not new_text:
            continue
            
        try:
            shape = prs.slides[slide_idx].shapes[shape_idx]
            if not shape.has_text_frame:
                continue
                
            # Activer l'auto-ajustement pour eviter les debordements visuels
            shape.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            shape.text_frame.word_wrap = True

            lines = new_text.split('\n')
            lines = [l for l in lines if l.strip()]

            # --- HARD TRUNCATION SAFEGUARD (V10) ---
            # Si l'IA n'ecoute pas et genere trop de texte pour une experience
            if slide_idx > 0:
                # V17: Retour a 100% de charset (les lignes eviteront l'overflow vertical toutes seules)
                max_chars = len(shape.text.strip())
                # V17: Hard limit sur le nombre de paragraphes/lignes (5, 5, 4 explicite)
                num_orig_paragraphs = len([p for p in shape.text_frame.paragraphs if p.text.strip()])
                if "AXA" in shape.text:
                    max_allowed_lines = 5
                elif "Helpline" in shape.text:
                    max_allowed_lines = 5
                elif "TSN" in shape.text:
                    max_allowed_lines = 3
                else:
                    max_allowed_lines = max(3, num_orig_paragraphs - 1)
                
                kept_lines = []
                current_chars = 0
                
                for line in lines:
                    # On est ultra strict : 0 tolerance char, ET limte stricte de nombre de puces
                    if len(kept_lines) < max_allowed_lines and current_chars + len(line) <= max_chars:
                        kept_lines.append(line)
                        current_chars += len(line) + 3
                    else:
                        print(f"  [WARN] Rognage line pour anti-overflow (Exp): {line[:30]}...")
                        # On arrete de garder des lignes pour cette experience
                        break
                        
                if not kept_lines and lines:
                    kept_lines = [lines[0][:max_chars]]
                lines = kept_lines

            num_new = len(lines)
            num_orig = len(shape.text_frame.paragraphs)
            
            # Sauvegarder les proprietes de base pour les nouveaux paragraphes
            font_name = None
            font_size = None
            font_bold = None
            font_color_rgb = None
            font_theme_color = None
            font_brightness = 0
            p_level = 0
            base_p_level = 0
            bullet_prefix = ""
            
            if num_orig > 0:
                p0 = shape.text_frame.paragraphs[0]
                p_level = p0.level
                
                # V12 : Extraire la puce manuelle ("-> " ou "* ") UNIQUEMENT pour la premiere ligne
                match = re.search(r"^([->*\->*]\s*)", p0.text)
                if match:
                    bullet_prefix = match.group(1)
                    
                # Le niveau de base pour les nouveaux paragraphes crees (pour continuer la liste)
                base_p_level = shape.text_frame.paragraphs[-1].level
                
                if p0.runs:
                    r0 = p0.runs[0]
                    font_name = r0.font.name
                    font_size = r0.font.size
                    font_bold = r0.font.bold
                    try:
                        if r0.font.color:
                            if r0.font.color.type == MSO_COLOR_TYPE.RGB:
                                font_color_rgb = r0.font.color.rgb
                            elif r0.font.color.type == MSO_COLOR_TYPE.THEME:
                                font_theme_color = r0.font.color.theme_color
                                font_brightness = r0.font.color.brightness
                    except: pass

            # Remplacement conservateur
            for i in range(max(num_new, num_orig)):
                if i < num_new:
                    if i < num_orig:
                        p = shape.text_frame.paragraphs[i]
                    else:
                        p = shape.text_frame.add_paragraph()
                        # V12 : Continuer la puce du dernier element pour les nouveaux ajouts
                        p.level = base_p_level
                    
                    final_line = lines[i].strip()
                    # 1. Nettoyage radical des puces manuelles generees par l'IA
                    # On enleve tout ce qui ressemble a une puce au debut (*, ->, -, *, >)
                    final_line = re.sub(r"^[->*\->*]\s*", "", final_line)
                    
                    # 2. V12 : Remettre la puce manuelle d'origine UNIQUEMENT sur la premiere ligne (i == 0)
                    if i == 0 and bullet_prefix:
                        prefix_clean = bullet_prefix.strip()
                        if not final_line.startswith(prefix_clean):
                            final_line = bullet_prefix + final_line
                            
                    p.text = final_line
                    
                    # Restaurer la mise en forme run par run
                    for r in p.runs:
                        if font_name: r.font.name = font_name
                        if font_size: r.font.size = font_size
                        if font_bold is not None: r.font.bold = font_bold
                        
                        # Fallback Noir : Si Slide 0 et block large (resume), ou si couleur non trouvee
                        # on force Noir (0,0,0) pour la lisibilite.
                        if slide_idx == 0 and sum(len(line) for line in lines) > 100:
                            # Pour le resume en haut, on force Noir
                            r.font.color.rgb = RGBColor(0, 0, 0)
                        elif font_color_rgb: 
                            r.font.color.rgb = font_color_rgb
                        elif font_theme_color is not None:
                            try:
                                r.font.color.theme_color = font_theme_color
                                r.font.color.brightness = font_brightness
                            except: pass
                else:
                    if i < len(shape.text_frame.paragraphs):
                        shape.text_frame.paragraphs[i].text = ""

        except Exception as e:
            print(f"  [WARN] Erreur lors de la mise a jour de la shape {shape_idx}: {e}")
            
    prs.save(output_path)

def convert_pptx_to_pdf(pptx_path, pdf_path):
    """Convertit un fichier PPTX en PDF via PowerPoint COM API."""
    powerpoint = None
    deck = None
    try:
        powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
        powerpoint.Visible = 1
        ppt_abs = os.path.abspath(pptx_path)
        pdf_abs = os.path.abspath(pdf_path)
        
        # S'assurer que le dossier parent du PDF existe
        os.makedirs(os.path.dirname(pdf_abs), exist_ok=True)
        
        deck = powerpoint.Presentations.Open(ppt_abs, WithWindow=False)
        # 32 est le format wdFormatPDF pour PowerPoint
        deck.SaveAs(pdf_abs, 32)
        return True
    except Exception as e:
        print(f"  [ERROR] Conversion PPTX vers PDF echouee: {e}")
        return False
    finally:
        if deck:
            deck.Close()
        if powerpoint:
            powerpoint.Quit()

def clean_json_response(text):
    """Nettoie la reponse de l'IA pour extraire uniquement le JSON valide."""
    text = text.strip()
    if "```json" in text:
        text = text.split("```json")[1].split("```")[0].strip()
    elif "```" in text:
        text = text.split("```")[1].split("```")[0].strip()
    return text

def adapter_cv_pptx(pptx_path, offre, dossier_sortie="cv modifi?"):
    """
    Lit le PPTX, identifie les blocs textes d'experience, 
    demande a l'IA de les adapter a l'offre, 
    genere un PPTX modifie, puis le convertit en PDF.
    """
    print(f"  ??-???  Adaptation du CV (PPTX) pour l'offre : {offre.get('titre', 'inconnu')}")
    
    shapes_data = extract_experience_shapes(pptx_path)
    if not shapes_data:
        print("  [WARN] Aucun bloc de texte exploitable trouve dans le PPTX.")
        return None
        
    entreprise = offre.get('entreprise', 'Non precise')
    titre_offre = offre.get('titre', 'Non precise')
    description_offre = offre.get('description', '')
    
    prompt = f"""
Tu es un expert en recrutement et optimisation ATS. Ton but est d'adapter les experiences d'un CV pour qu'elles "matchent" parfaitement l'offre d'emploi ciblee.

CRITERE DE REUSSITE : Ton texte doit maximiser le score ATS en integrant naturellement les mots-cles de l'offre.

Offre ciblee :
Poste : {titre_offre}
Entreprise : {entreprise}
Description :
{description_offre[:2500]}

Voici les differents blocs de texte extraits du CV actuel (format JSON) :
{json.dumps(shapes_data, ensure_ascii=False, indent=2)}

Instructions :
1. Identifie uniquement les blocs d'EXPERIENCES PROFESSIONNELLES, de COMPETENCES ou le RESUME/PROFIL (au debut). 
2. **ZERO HALLUCINATION** : N'invente JAMAIS d'annees d'experience. Si le CV n'en mentionne pas 5, n'en invente pas 5. Utilise les donnees reelles.
3. **CIBLE DE LONGUEUR ET DE PUCES (EXTR?MEMENT CRITIQUE - NE PAS D?BORDER)** : 
   - Pour CHAQUE bloc, je t'ai fourni `min_chars`, `max_chars` ET surtout `exact_bullet_count` dans le JSON.
   - Ton nouveau texte DOIT OBLIGATOIREMENT faire entre `min_chars` et `max_chars` caract?res au total.
   - Tu DOIS OBLIGATOIREMENT g?n?rer EXACTEMENT `exact_bullet_count` lignes (bullet points) pour ce bloc. Ni plus, ni moins.
   - Ex: Si max_chars = 661 et exact_bullet_count = 5, ton texte ne DOIT JAMAIS FAIRE 662 caract?res, et DOIT AVOIR EXACTEMENT 5 phrases s?par?es par de simples retours ? la ligne `\\n`.
   - L'exp?rience Helpline et TSN DOIVENT etre denses (proche de max_chars), r?partis exactement sur le nombre de puces demand?.
4. **STRUCTURE DU RESUME ( Slide 0 )** : 
   - Ne mets AUCUN texte en gras. N'utilise pas les etoiles Markdown.
   - Formate exactement en 3 lignes/paragraphes separes par `\\n` dans le JSON :
     Ligne 1 : Pr?sentation courte.
     Ligne 2 : Comp?tences : [Liste concise].
     Ligne 3 : Valeur ajout?e : [Phrase courte].
5. **STRUCTURE DES EXPERIENCES (CRITIQUE)** :
   - Fais des BULLET POINTS uniquement.
   - SEPARE tes points par des sauts de ligne (`\\n` dans le JSON).
   - N'ajoute AUCUNE puce (*, ->, -) par toi-meme au debut de tes lignes. Le script PowerPoint s'en charge. Commence directement par la premiere lettre du point.
6. Renvoie UNIQUEMENT un tableau JSON. Pas de texte avant ou apres.

Format de sortie attendu :
[
  {{
    "slide_idx": 1,
    "shape_idx": 4,
    "text": "Nouvelle description bullet point 1\\nNouvelle description bullet point 2"
  }}
]
"""
    
    response_text = _call_openai_with_mistral_fallback(prompt, max_tokens=2000)
    json_str = clean_json_response(response_text)
    
    try:
        updates = json.loads(json_str)
    except json.JSONDecodeError as e:
        print(f"  [ERROR] L'IA n'a pas renvoye un JSON valide : {e}")
        return None
        
    if not isinstance(updates, list) or len(updates) == 0:
        print("  [INFO] L'IA n'a propose aucune modification pour ce CV.")
        return None
        
    # Creer des noms de fichiers clairs
    # On nettoie le nom de l'entreprise et du poste pour le nom du fichier
    safe_entreprise = re.sub(r'[\\\\/*?:"<>|]', "", entreprise).strip()
    safe_titre = re.sub(r'[\\\\/*?:"<>|]', "", titre_offre).strip()
    safe_entreprise = safe_entreprise[:30] # Limiter la longueur
    safe_titre = safe_titre[:30]
    
    dossier_sortie_abs = os.path.abspath(dossier_sortie)
    os.makedirs(dossier_sortie_abs, exist_ok=True)
    
    # Fichier PPTX temp
    temp_pptx = os.path.join(dossier_sortie_abs, "temp_cv.pptx")
    update_pptx_shapes(pptx_path, updates, temp_pptx)
    
    # PDF a uploader (nom dynamique base sur config)
    annee_actuelle = datetime.datetime.now().year
    nom_cv = f"{config.PRENOM} {config.NOM.upper()} {annee_actuelle}.pdf"
    upload_pdf_path = os.path.abspath(nom_cv)
    
    # PDF d'archive pour le suivi (nom dynamique complet)
    archive_pdf_name = f"{config.PRENOM} {config.NOM.upper()} - {safe_entreprise} - {safe_titre}.pdf"
    archive_pdf_path = os.path.join(dossier_sortie_abs, archive_pdf_name)
    
    print("   [WAIT] Conversion du PPTX modifie en PDF...")
    success = convert_pptx_to_pdf(temp_pptx, archive_pdf_path)
    
    if success:
        import shutil
        # Copier l'archive vers le fichier a uploader pour avoir le bon nom lors de l'upload
        shutil.copy2(archive_pdf_path, upload_pdf_path)
        
        # Nettoyer le PPTX temp
        try:
            os.remove(temp_pptx)
        except Exception:
            pass
            
        print(f"   [OK] CV personnalise genere et converti: {archive_pdf_path}")
        return upload_pdf_path
    else:
        return None
