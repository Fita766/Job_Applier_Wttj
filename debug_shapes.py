import sys
from pptx import Presentation

prs = Presentation(sys.argv[1])
slide = prs.slides[0]

for i, shape in enumerate(slide.shapes):
    if shape.has_text_frame:
        if "Zoho" in shape.text:
            print(f"Shape {i}:")
            print(shape.text)
            print("-" * 20)
