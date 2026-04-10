import sys
from pptx import Presentation

def analyze_pptx(pptx_path):
    with open("pptx_structure.txt", "w", encoding="utf-8") as f:
        f.write(f"Analyzing {pptx_path}...\n")
        prs = Presentation(pptx_path)
        for i, slide in enumerate(prs.slides):
            f.write(f"\n--- Slide {i+1} ---\n")
            for j, shape in enumerate(slide.shapes):
                if not shape.has_text_frame:
                    continue
                text = shape.text.strip()
                if not text:
                    continue
                
                # Print first 100 chars
                preview = text.replace('\\n', ' ')[:100]
                f.write(f"  Shape {j}: {preview}...\n")
                
                # Print runs for the first few paragraphs to see formatting
                for p_idx, p in enumerate(shape.text_frame.paragraphs):
                    if p_idx > 2:
                        break
                    runs_info = [f"'{r.text}'" for r in p.runs]
                    if runs_info:
                        f.write(f"    P{p_idx} runs: {', '.join(runs_info)}\n")

if __name__ == "__main__":
    analyze_pptx(sys.argv[1])

